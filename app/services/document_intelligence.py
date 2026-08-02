"""
Stage 1: Document Intelligence.

Parses an uploaded file into a ParsedDocument, preserving structural signal
(headings, tables, figures) rather than collapsing everything to flat text.
This matters because Stage 2 (classification) and Stage 3 (knowledge
extraction) both work much better with structure-aware input than with a
single undifferentiated text blob  a heading tells the model "this is a
section boundary", which a wall of text doesn't.

Supported formats: PDF, DOCX, PPTX, plain text. PDFs additionally get
routed between a direct-text-extraction path and an OCR path, since NCERT
chapters (the stated reference document type) are frequently scanned pages
rather than text-layer PDFs, and pdfplumber silently returns empty or
near-empty text on a scanned page rather than raising an error -- without
this routing, a scanned chapter would flow through the rest of the
pipeline as an almost-blank document with no warning.
"""

import io
import re
from pathlib import Path
from collections.abc import Callable

from app.models.schemas import DocumentBlock, DocumentTypeHint, ParsedDocument

# Below this average extracted characters per page, a PDF is treated as
# likely-scanned even if the user picked "not sure" or guessed wrong --
# a real text-layer page of a textbook chapter is rarely this sparse.
_SCANNED_HEURISTIC_CHARS_PER_PAGE = 40


def parse_document(
    filename: str,
    file_bytes: bytes,
    doc_type_hint: DocumentTypeHint = DocumentTypeHint.NOT_SURE,
    progress_callback: Callable[[int], None] | None = None,
) -> ParsedDocument:
    ext = Path(filename).suffix.lower()

    if progress_callback is not None:
        progress_callback(0)

    if ext == ".pdf":
        blocks, meta = _parse_pdf_routed(file_bytes, doc_type_hint, progress_callback)
    elif ext == ".docx":
        blocks, meta = _parse_docx(file_bytes, progress_callback)
    elif ext in (".pptx", ".ppt"):
        blocks, meta = _parse_pptx(file_bytes, progress_callback)
    elif ext in (".txt", ".md"):
        blocks, meta = _parse_text(file_bytes, progress_callback)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if progress_callback is not None:
        progress_callback(100)

    raw_text = "\n".join(b.text for b in blocks)
    return ParsedDocument(
        source_filename=filename,
        raw_text=raw_text,
        blocks=blocks,
        metadata=meta,
    )


def _parse_pdf_routed(
    file_bytes: bytes,
    doc_type_hint: DocumentTypeHint,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    """
    Cost-aware routing per the client's stated preference: a user hint
    picks the parsing path directly where it's unambiguous (mostly_text /
    text_with_tables / text_with_diagrams / text_with_equations all mean
    "this PDF has a real text layer, use the cheap path"; scanned_pdf means
    "go straight to OCR"). When the hint is not_sure, fall back to trying
    the cheap text-extraction path first and only pay for OCR if that
    heuristically looks like it failed (very little text per page), rather
    than OCR'ing every "not sure" upload by default.
    """
    if doc_type_hint == DocumentTypeHint.SCANNED_PDF:
        return _parse_pdf_ocr(file_bytes, progress_callback)

    blocks, meta = _parse_pdf_text(file_bytes, progress_callback)

    if doc_type_hint == DocumentTypeHint.NOT_SURE:
        page_count = meta.get("page_count", 1) or 1
        total_chars = sum(len(b.text) for b in blocks)
        if total_chars / page_count < _SCANNED_HEURISTIC_CHARS_PER_PAGE:
            ocr_blocks, ocr_meta = _parse_pdf_ocr(file_bytes, progress_callback)
            ocr_meta["routing"] = "auto-detected-scanned (text-layer extraction was too sparse)"
            return ocr_blocks, ocr_meta

    meta["routing"] = f"text-layer extraction (hint: {doc_type_hint.value})"
    return blocks, meta


def _parse_pdf_ocr(
    file_bytes: bytes,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    """
    Rasterizes each page and runs Tesseract OCR over it. Requires the
    `tesseract-ocr` system package to be installed on the host (e.g.
    `apt-get install tesseract-ocr` on the deployment image) -- pytesseract
    is just a wrapper around that binary, not a pure-Python OCR engine.
    """
    import fitz  # PyMuPDF
    import pytesseract
    from PIL import Image

    blocks: list[DocumentBlock] = []
    pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
    page_count = len(pdf_doc)

    try:
        for page_num, page in enumerate(pdf_doc, start=1):
            if progress_callback is not None and page_count:
                progress_callback(max(0, min(99, int((page_num - 1) / page_count * 100))))

            # 2x zoom roughly doubles effective DPI over the PDF's base
            # rendering resolution, which noticeably improves OCR accuracy on
            # textbook-scan-quality source images without a large runtime cost.
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            page_text = pytesseract.image_to_string(image)

            for line in page_text.split("\n"):
                stripped = line.strip()
                if not stripped:
                    continue
                block_type, level = _classify_line(stripped)
                blocks.append(
                    DocumentBlock(block_type=block_type, level=level, text=stripped, page=page_num)
                )

            if progress_callback is not None and page_count:
                progress_callback(max(1, min(99, int(page_num / page_count * 100))))
    except pytesseract.TesseractNotFoundError:
        # Graceful fallback: if the host doesn't have Tesseract installed,
        # don't kill the whole pipeline during document parsing. Fall back to
        # direct text extraction so text-layer PDFs still work, and scanned
        # PDFs produce a clearly degraded result instead of a hard crash.
        pdf_doc.close()
        blocks, meta = _parse_pdf_text(file_bytes, progress_callback)
        meta["parsed_via"] = "text-layer fallback (tesseract unavailable)"
        meta["warning"] = (
            "Tesseract OCR is not installed or not on PATH, so scanned PDFs will be "
            "parsed without OCR. Install tesseract-ocr for best results."
        )
        return blocks, meta

    pdf_doc.close()
    return blocks, {"format": "pdf", "page_count": page_count, "parsed_via": "ocr"}


def _parse_pdf_text(
    file_bytes: bytes,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    import pdfplumber

    blocks: list[DocumentBlock] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        page_count = len(pdf.pages)
        for page_num, page in enumerate(pdf.pages, start=1):
            if progress_callback is not None and page_count:
                progress_callback(max(0, min(99, int((page_num - 1) / page_count * 100))))

            # Tables first, since plain text extraction on a page with a
            # table often mangles column alignment.
            tables = page.extract_tables()
            for table in tables:
                table_text = "\n".join(
                    " | ".join(cell or "" for cell in row) for row in table
                )
                blocks.append(
                    DocumentBlock(block_type="table", text=table_text, page=page_num)
                )

            text = page.extract_text() or ""
            for line in text.split("\n"):
                stripped = line.strip()
                if not stripped:
                    continue
                block_type, level = _classify_line(stripped)
                blocks.append(
                    DocumentBlock(
                        block_type=block_type,
                        level=level,
                        text=stripped,
                        page=page_num,
                    )
                )

            if progress_callback is not None and page_count:
                progress_callback(max(1, min(99, int(page_num / page_count * 100))))

    return blocks, {"format": "pdf", "page_count": page_count}


def _parse_docx(
    file_bytes: bytes,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    import docx

    document = docx.Document(io.BytesIO(file_bytes))
    blocks: list[DocumentBlock] = []

    if progress_callback is not None:
        progress_callback(20)

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if style.startswith("heading"):
            level = int(style.replace("heading", "").strip() or 1)
            blocks.append(DocumentBlock(block_type="heading", level=level, text=text))
        else:
            blocks.append(DocumentBlock(block_type="paragraph", text=text))

    for table in document.tables:
        table_text = "\n".join(
            " | ".join(cell.text for cell in row.cells) for row in table.rows
        )
        blocks.append(DocumentBlock(block_type="table", text=table_text))

    if progress_callback is not None:
        progress_callback(100)

    return blocks, {"format": "docx"}


def _parse_pptx(
    file_bytes: bytes,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    blocks: list[DocumentBlock] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if progress_callback is not None and len(prs.slides):
            progress_callback(max(0, min(99, int((slide_num - 1) / len(prs.slides) * 100))))

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                # Title placeholders act as slide-level headings
                is_title = getattr(shape, "is_placeholder", False) and getattr(
                    shape.placeholder_format, "type", None
                ) is not None and shape == slide.shapes.title
                blocks.append(
                    DocumentBlock(
                        block_type="heading" if is_title else "paragraph",
                        level=1 if is_title else None,
                        text=text,
                        page=slide_num,
                    )
                )

        if progress_callback is not None and len(prs.slides):
            progress_callback(max(1, min(99, int(slide_num / len(prs.slides) * 100))))

    return blocks, {"format": "pptx", "slide_count": len(prs.slides)}


def _parse_text(
    file_bytes: bytes,
    progress_callback: Callable[[int], None] | None = None,
) -> tuple[list[DocumentBlock], dict]:
    text = file_bytes.decode("utf-8", errors="ignore")
    if progress_callback is not None:
        progress_callback(100)

    blocks: list[DocumentBlock] = []
    lines = text.split("\n")
    in_code_fence = False
    code_lines: list[str] = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code_fence:
                # Closing fence -- flush the accumulated code as one block
                # rather than classifying each line separately, since a
                # line-by-line pass over source code sends most lines
                # through the same generic paragraph/equation heuristics
                # meant for prose, with poor results (see _classify_line).
                if code_lines:
                    blocks.append(DocumentBlock(block_type="code", text="\n".join(code_lines)))
                code_lines = []
                in_code_fence = False
            else:
                in_code_fence = True
            continue

        if in_code_fence:
            code_lines.append(line)
            continue

        if not stripped:
            continue
        block_type, level = _classify_line(stripped)
        blocks.append(DocumentBlock(block_type=block_type, level=level, text=stripped))

    if code_lines:  # unterminated fence at end of file -- flush what we have
        blocks.append(DocumentBlock(block_type="code", text="\n".join(code_lines)))

    return blocks, {"format": "text"}


def _classify_line(line: str) -> tuple[str, int | None]:
    """
    Cheap structural heuristics for plain-text and PDF-extracted lines,
    where there is no real style metadata to lean on the way DOCX gives us.
    Kept intentionally simple: this is a signal for the LLM stages to lean
    on, not a ground-truth structure parser.

    Markdown "#"-prefixed lines are recognized explicitly as headings
    (stripped of the leading hashes) before anything else is checked --
    without this, real section headings in a markdown-formatted document
    fall through to the generic heuristics below and get misclassified as
    plain paragraphs.
    """
    markdown_heading = re.match(r"^(#{1,6})\s+(.+)$", line)
    if markdown_heading:
        return "heading", len(markdown_heading.group(1))

    # Equation detection requires an actual math symbol, not just an
    # equals sign. A bare "=" matches nearly every line of source code
    # (assignment statements), which previously meant technical documents
    # with embedded code got most of their code lines mislabeled as
    # equations rather than paragraphs/code -- verified against a real
    # Python-heavy markdown document where this fired on >200 lines like
    # `s1 = pd.Series([1, 2, 3, 4])`. Genuine math notation is a much
    # narrower, more reliable signal here.
    if any(sym in line for sym in ("\u2211", "\u222b", "\u2202", "\u221a", "\u00b1", "\u2264", "\u2265", "\u2260", "\u0394")) and len(line) < 120:
        return "equation", None

    if line.isupper() and len(line.split()) <= 12:
        return "heading", 1

    # A numbered line that's actually a markdown link (table of contents
    # entries like "1. [Pandas Fundamentals](#1-pandas-fundamentals)") is
    # a reference to a heading, not the heading itself -- skip it here so
    # it doesn't get misclassified ahead of the real section heading it
    # points to.
    if line[0].isdigit() and "." in line[:4] and len(line) < 100 and "](" not in line:
        return "heading", 2

    return "paragraph", None
